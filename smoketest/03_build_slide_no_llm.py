#!/usr/bin/env python3
"""Stage 3 — build a real .pptx from a bundled template with NO LLM at all.

This is the load-bearing test. It walks the exact code path the MCP server uses
for generate_slide()/save_generated_slides(), except that the edit actions are
hand-written here instead of produced by the `coder` LLM:

    deepcopy(template slide)
      -> CodeExecutor.execute_actions(actions, slide, doc)   # pptagent/apis.py:127
      -> Presentation.validate(slide)                        # drops untouched paragraphs
      -> Presentation.save(path)                             # python-pptx write

If this passes, everything except the model call is proven to work offline:
no LibreOffice, no Docker, no network.

Run:  python smoketest/03_build_slide_no_llm.py [--template default] [--slides 3]
"""

from __future__ import annotations

import argparse
import json
import sys
from copy import deepcopy
from pathlib import Path

from pptagent.apis import CodeExecutor
from pptagent.multimodal import ImageLabler
from pptagent.presentation import Presentation
from pptagent.presentation.layout import Layout
from pptagent.utils import Config, Language, package_join

HERE = Path(__file__).resolve().parent
OUT_JSON = HERE / "_result_03_build.json"

# 슬라이드에 채워 넣을 한국어 내용. 실제로는 Research 단계 원고에서 온다.
CONTENT: list[list[str]] = [
    ["사내 PPT 자동 생성 검토", "템플릿 기반 파이프라인 오프라인 실행 테스트", "2026-08-18"],
    ["검토 배경", "LibreOffice 없이도 슬라이드 생성이 가능한지 확인한다",
     "python-pptx만으로 템플릿 편집이 끝나는지 본다", "LLM 연결은 다음 단계로 미룬다"],
    ["확인 결과", "번들 템플릿 6종은 모두 오프라인 로드된다",
     "슬라이드 편집과 저장은 모델 없이 동작한다", "모델은 편집 코드 생성에만 쓰인다"],
]


def load_template(name: str):
    """Load a template the same way PPTAgentServer.__init__ does."""
    folder = Path(package_join("templates")) / name
    cfg = Config(str(folder))
    prs = Presentation.from_file(str(folder / "source.pptx"), cfg)
    ImageLabler(prs, cfg).apply_stats(json.loads((folder / "image_stats.json").read_text()))

    induction = dict(json.loads((folder / "slide_induction.json").read_text()))
    lang = Language(**induction.pop("language"))
    functional = induction.pop("functional_keys")
    layouts = {k: Layout(title=k, **v) for k, v in induction.items()}
    return prs, layouts, functional, lang


def discover_paragraphs(slide) -> list[tuple[int, int, str]]:
    """List every addressable (div_id, paragraph_id, current text) on a slide.

    div_id is shape.shape_idx and paragraph_id is para.idx — exactly the two
    integers replace_paragraph()/del_paragraph() take (pptagent/apis.py:357,403).
    """
    found: list[tuple[int, int, str]] = []
    for shape in slide:
        if not shape.text_frame.is_textframe:
            continue
        for para in shape.text_frame.paragraphs:
            if para.idx == -1 or not para.text.strip():
                continue
            found.append((shape.shape_idx, para.idx, para.text))
    return found


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", default="default")
    ap.add_argument("--slides", type=int, default=3)
    ap.add_argument("--out", default=str(HERE / "smoketest_output.pptx"))
    args = ap.parse_args()

    print(f"\n=== Stage 3: LLM 없이 슬라이드 생성 (template={args.template}) ===\n")

    prs, layouts, functional, lang = load_template(args.template)
    empty_prs = deepcopy(prs)  # set_reference()가 만드는 빈 사본과 같은 역할
    print(f"템플릿 로드: 슬라이드 {len(prs.slides)}장, 레이아웃 {len(layouts)}개, 언어 {lang.lid}")

    # 텍스트 위주 레이아웃을 우선 고른다 (이미지 교체는 Document 객체가 필요해서 제외)
    ordered = sorted(
        layouts.values(),
        key=lambda lo: (any(el.type == "image" for el in lo.elements), lo.title != "opening"),
    )
    chosen = ordered[: args.slides]
    print("선택한 레이아웃:")
    for lo in chosen:
        print(f"  · {lo.title}  (template_id={lo.template_id})")

    built = []
    report: list[dict] = []

    for i, layout in enumerate(chosen):
        content = CONTENT[i % len(CONTENT)]
        print(f"\n--- 슬라이드 {i + 1}: {layout.title[:50]} ---")

        edit_slide = deepcopy(prs.slides[layout.template_id - 1])
        paras = discover_paragraphs(edit_slide)
        print(f"편집 가능한 문단 {len(paras)}개:")
        for div_id, para_id, text in paras[:8]:
            print(f"    div={div_id:<3} para={para_id:<3} {text[:46]!r}")

        # coder LLM이 만들어야 할 편집 코드를 여기서는 직접 작성한다.
        lines = []
        for (div_id, para_id, _), new_text in zip(paras, content):
            lines.append(f'replace_paragraph({div_id}, {para_id}, "{new_text}")')
        actions = "\n".join(lines)
        if not actions:
            print("  [FAIL] 편집 가능한 문단이 없어 건너뜀")
            report.append({"layout": layout.title, "ok": False, "reason": "no paragraph"})
            continue

        print("생성한 편집 액션:")
        for ln in lines:
            print(f"    {ln}")

        executor = CodeExecutor(retry_times=1)
        executor.command_history.append([None, "smoketest", None])
        feedback = executor.execute_actions(actions, edit_slide, None)
        if feedback is not None:
            print(f"  [FAIL] 실행 오류:\n{feedback[1][-500:]}")
            report.append({"layout": layout.title, "ok": False, "reason": feedback[1][-300:]})
            continue

        # 손대지 않은 문단은 여기서 삭제 대상으로 표시된다
        empty_prs.validate(edit_slide)
        built.append(edit_slide)
        print(f"  [OK  ] {len(lines)}개 문단 교체 완료")
        report.append({"layout": layout.title, "ok": True, "edits": len(lines)})

    if not built:
        print("\n생성된 슬라이드가 없습니다.\n")
        return 1

    out_path = Path(args.out)
    empty_prs.slides = built
    empty_prs.save(str(out_path))
    size = out_path.stat().st_size
    print(f"\n저장: {out_path}  ({size:,} bytes, {len(built)}장)")

    # ── 독립 검증: python-pptx로 다시 열어 텍스트를 확인한다 ──────────
    print("\n--- python-pptx 독립 검증 ---")
    from pptx import Presentation as PptxPresentation

    check = PptxPresentation(str(out_path))
    print(f"  슬라이드 수: {len(check.slides)}")
    texts_found = 0
    for n, slide in enumerate(check.slides, 1):
        got = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                got.append(shape.text_frame.text.strip().replace("\n", " / ")[:70])
        texts_found += len(got)
        print(f"  슬라이드 {n}:")
        for t in got:
            print(f"      {t}")

    expected_texts = {t for row in CONTENT[: len(built)] for t in row}
    all_text = " ".join(
        shape.text_frame.text
        for slide in check.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    hit = sum(1 for t in expected_texts if t in all_text)
    print(f"\n  삽입한 문자열 {len(expected_texts)}개 중 {hit}개가 파일에서 확인됨")

    OUT_JSON.write_text(
        json.dumps(
            {"template": args.template, "slides": report, "output": str(out_path),
             "bytes": size, "verified_strings": f"{hit}/{len(expected_texts)}"},
            ensure_ascii=False, indent=2,
        ),
        encoding="utf-8",
    )
    ok = len(check.slides) == len(built) and hit >= len(expected_texts) * 0.6
    print(f"\n  [{'OK  ' if ok else 'FAIL'}] Stage 3 {'통과' if ok else '실패'}")
    print(f"  결과 저장: {OUT_JSON}\n")
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
